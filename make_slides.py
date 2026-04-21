"""
make_slides.py — Generate presentation slides for Finance 372T Group 9 demo.
Run: python make_slides.py
Output: presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Brand colors
# ---------------------------------------------------------------------------
BURNT_ORANGE = RGBColor(0x8B, 0x3A, 0x00)
DARK_GRAY    = RGBColor(0x33, 0x33, 0x33)
MID_GRAY     = RGBColor(0x77, 0x77, 0x77)
LIGHT_GRAY   = RGBColor(0xF4, 0xF4, 0xF4)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_ORANGE = RGBColor(0xFF, 0xCC, 0xAA)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def rect(slide, l, t, w, h, fill=None, line_color=None, line_w=Pt(0.75)):
    s = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_color:
        s.line.color.rgb = line_color; s.line.width = line_w
    else:
        s.line.fill.background()
    return s


def txt(slide, text, l, t, w, h,
        size=18, bold=False, italic=False,
        color=DARK_GRAY, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return box


def bullets(slide, items, l, t, w, h, size=16, color=DARK_GRAY, gap=Pt(10)):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = gap
        r = p.add_run()
        r.text = f"\u2022  {item}"
        r.font.size = Pt(size); r.font.color.rgb = color
    return box


def header(slide, title, subtitle=None):
    rect(slide, 0, 0, SLIDE_W, Inches(1.05), fill=BURNT_ORANGE)
    txt(slide, title,
        Inches(0.4), Inches(0.1), Inches(12), Inches(0.6),
        size=30, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle,
            Inches(0.4), Inches(0.65), Inches(12), Inches(0.32),
            size=12, color=LIGHT_ORANGE)


def footer(slide):
    rect(slide, 0, Inches(7.18), SLIDE_W, Inches(0.32), fill=DARK_GRAY)
    txt(slide, "Finance 372T/397  \u00b7  Group 9  \u00b7  Blake Stanley & Will Pechersky  \u00b7  April 21, 2026",
        Inches(0.3), Inches(7.2), Inches(12.5), Inches(0.25),
        size=10, color=WHITE)


def divider(slide, t, l=Inches(0.4), w=Inches(12.5)):
    rect(slide, l, t, w, Pt(1.5), fill=BURNT_ORANGE)


def numbered_circle(slide, num, l, t, size=Inches(0.45)):
    c = slide.shapes.add_shape(9, l, t, size, size)
    c.fill.solid(); c.fill.fore_color.rgb = BURNT_ORANGE
    c.line.fill.background()
    tf = c.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = str(num); r.font.size = Pt(14)
    r.font.bold = True; r.font.color.rgb = WHITE


# ---------------------------------------------------------------------------
# SLIDE 1 — The Problem  (what & why)
# ---------------------------------------------------------------------------

s1 = prs.slides.add_slide(blank)
rect(s1, 0, 0, SLIDE_W, SLIDE_H, fill=LIGHT_GRAY)
header(s1, "The Problem", "Why equity research needs a better tool")
footer(s1)

# Large central statement
rect(s1, Inches(1.5), Inches(1.3), Inches(10.3), Inches(1.8), fill=WHITE,
     line_color=RGBColor(0xDD, 0xDD, 0xDD))
txt(s1, "\u201cInvestors miss divergences between what the",
    Inches(1.8), Inches(1.45), Inches(9.8), Inches(0.6),
    size=22, bold=True, color=BURNT_ORANGE, align=PP_ALIGN.CENTER)
txt(s1, "numbers say and what management says.\u201d",
    Inches(1.8), Inches(1.9), Inches(9.8), Inches(0.5),
    size=22, bold=True, color=BURNT_ORANGE, align=PP_ALIGN.CENTER)

# 3 problem cards
card_w = Inches(3.8)
card_h = Inches(2.8)
card_top = Inches(3.3)
problems = [
    ("Scattered Data", "Fundamentals, filings, and sentiment signals live in separate terminals and databases"),
    ("No Unified View", "Free tools don\u2019t combine quant signals with earnings call analysis in one place"),
    ("Missed Signal", "A stock can look strong on paper while management tone is quietly deteriorating"),
]
for i, (title, body) in enumerate(problems):
    l = Inches(0.4) + i * (card_w + Inches(0.2))
    rect(s1, l, card_top, card_w, card_h, fill=WHITE,
         line_color=RGBColor(0xDD, 0xDD, 0xDD))
    rect(s1, l, card_top, card_w, Inches(0.06), fill=BURNT_ORANGE)
    txt(s1, title, l + Inches(0.15), card_top + Inches(0.15), card_w - Inches(0.3), Inches(0.4),
        size=16, bold=True, color=BURNT_ORANGE)
    txt(s1, body, l + Inches(0.15), card_top + Inches(0.6), card_w - Inches(0.3), Inches(2.0),
        size=14, color=DARK_GRAY)


# ---------------------------------------------------------------------------
# SLIDE 2 — The Solution  (what & why continued + MVP)
# ---------------------------------------------------------------------------

s2 = prs.slides.add_slide(blank)
rect(s2, 0, 0, SLIDE_W, SLIDE_H, fill=LIGHT_GRAY)
header(s2, "The Solution", "Input any U.S. equity ticker \u2192 instant research dashboard")
footer(s2)

txt(s2, "[ Dashboard screenshot \u2014 paste here ]",
    Inches(0.4), Inches(1.2), Inches(7.6), Inches(5.7),
    size=16, color=MID_GRAY, align=PP_ALIGN.CENTER)
rect(s2, Inches(0.4), Inches(1.2), Inches(7.6), Inches(5.7),
     line_color=RGBColor(0xCC, 0xCC, 0xCC))

# Module list on the right
modules = [
    ("1", "Piotroski F-Score",      "Financial strength \u00b7 9 components"),
    ("2", "Gross Profitability",     "Novy-Marx signal \u00b7 sector percentile"),
    ("3", "Earnings Quality",        "Accruals ratio \u00b7 cash vs. net income"),
    ("4", "Valuation & Momentum",   "EV/EBITDA \u00b7 P/E \u00b7 12-1M momentum"),
    ("5", "Earnings Call Sentiment", "Tone \u00b7 hedging \u00b7 QoQ trend"),
    ("6", "AI Synthesis",            "Claude \u2192 bull / bear / risks / divergence"),
]
for i, (num, title, desc) in enumerate(modules):
    top = Inches(1.3) + i * Inches(0.88)
    numbered_circle(s2, num, Inches(8.2), top + Inches(0.04))
    txt(s2, title, Inches(8.75), top, Inches(4.3), Inches(0.35),
        size=14, bold=True, color=DARK_GRAY)
    txt(s2, desc, Inches(8.75), top + Inches(0.35), Inches(4.3), Inches(0.35),
        size=12, color=MID_GRAY)


# ---------------------------------------------------------------------------
# SLIDE 3 — Live Demo
# ---------------------------------------------------------------------------

s3 = prs.slides.add_slide(blank)
rect(s3, 0, 0, SLIDE_W, SLIDE_H, fill=BURNT_ORANGE)
footer(s3)

txt(s3, "Live Demo", Inches(0), Inches(2.5), SLIDE_W, Inches(1.2),
    size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s3, "AAPL  \u00b7  MSFT  \u00b7  HOG",
    Inches(0), Inches(3.8), SLIDE_W, Inches(0.6),
    size=24, color=LIGHT_ORANGE, align=PP_ALIGN.CENTER)
txt(s3, "streamlit run dashboard/app.py",
    Inches(0), Inches(4.6), SLIDE_W, Inches(0.5),
    size=16, italic=True, color=RGBColor(0xFF, 0xEE, 0xDD),
    align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# SLIDE 4 — How It Works  (methodology)
# ---------------------------------------------------------------------------

s4 = prs.slides.add_slide(blank)
rect(s4, 0, 0, SLIDE_W, SLIDE_H, fill=LIGHT_GRAY)
header(s4, "How It Works", "Architecture \u00b7 Data \u00b7 AI")
footer(s4)

# Architecture flow: 4 boxes with arrows
boxes = [
    ("Ticker Input",         "User types any\nU.S. equity symbol"),
    ("Quant Signals",        "WRDS Compustat + CRSP\nPre-computed Dec 2024\n33,675 tickers\n5 signals \u2192 composite score"),
    ("Sentiment Pipeline",   "Motley Fool (live)\nVADER tone score\nL-M hedging score\nQoQ trend tracker"),
    ("Claude Synthesis",     "Sonnet 4.6\nPrompt caching\nBull / Bear / Risks\nDivergence flag"),
]
box_w = Inches(2.8)
box_h = Inches(3.8)
box_top = Inches(1.8)
arrow_w = Inches(0.5)
total = len(boxes) * box_w + (len(boxes) - 1) * arrow_w
start_l = (SLIDE_W - total) / 2

for i, (title, body) in enumerate(boxes):
    l = start_l + i * (box_w + arrow_w)
    rect(s4, l, box_top, box_w, box_h, fill=WHITE,
         line_color=RGBColor(0xDD, 0xDD, 0xDD))
    rect(s4, l, box_top, box_w, Inches(0.45), fill=BURNT_ORANGE)
    txt(s4, title, l + Inches(0.1), box_top + Inches(0.05),
        box_w - Inches(0.2), Inches(0.35),
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s4, body, l + Inches(0.15), box_top + Inches(0.6),
        box_w - Inches(0.3), Inches(3.0),
        size=13, color=DARK_GRAY)
    # Arrow (except after last box)
    if i < len(boxes) - 1:
        arrow_l = l + box_w
        txt(s4, "\u2192", arrow_l, box_top + Inches(1.6), arrow_w, Inches(0.5),
            size=28, bold=True, color=BURNT_ORANGE, align=PP_ALIGN.CENTER)

# Bottom note on stale data
rect(s4, Inches(1.5), Inches(5.85), Inches(10.3), Inches(0.7),
     fill=WHITE, line_color=RGBColor(0xDD, 0xCC, 0xAA))
txt(s4, "\u26a0  Quant signals reflect Dec 2024 data (pre-computed from WRDS). "
        "Sentiment is fetched live. Stale data is an acknowledged limitation \u2014 "
        "appropriate for an academic research tool.",
    Inches(1.7), Inches(5.9), Inches(9.9), Inches(0.6),
    size=12, color=MID_GRAY, italic=True)


# ---------------------------------------------------------------------------
# SLIDE 5 — Results & Takeaways  (project quality)
# ---------------------------------------------------------------------------

s5 = prs.slides.add_slide(blank)
rect(s5, 0, 0, SLIDE_W, SLIDE_H, fill=LIGHT_GRAY)
header(s5, "Results & Takeaways")
footer(s5)

# Left: example outputs
rect(s5, Inches(0.3), Inches(1.2), Inches(8.2), Inches(5.7),
     fill=WHITE, line_color=RGBColor(0xDD, 0xDD, 0xDD))
txt(s5, "EXAMPLE OUTPUTS", Inches(0.5), Inches(1.3), Inches(4), Inches(0.3),
    size=10, bold=True, color=BURNT_ORANGE)
divider(s5, Inches(1.6), l=Inches(0.5), w=Inches(7.8))

examples = [
    ("AAPL", "Composite 62/100 \u00b7 F-Score 6/9 \u00b7 Strong gross profitability (84th pct)\n"
              "Expensive valuation (EV/EBITDA 87th pct) \u00b7 Tone mildly positive, low hedging\n"
              "Claude: durable moat offset by premium valuation with little margin of safety"),
    ("MSFT", "Strong fundamentals across all 5 signals \u00b7 Sentiment improving QoQ\n"
              "No divergence flagged \u2014 quant and management tone aligned"),
    ("HOG",  "Mid-cap contrast \u00b7 Lower composite score \u00b7 Different sector story\n"
              "Demonstrates dashboard works across cap ranges"),
]
top = Inches(1.75)
for ticker, body in examples:
    rect(s5, Inches(0.5), top, Inches(7.8), Inches(1.45),
         fill=LIGHT_GRAY, line_color=RGBColor(0xCC, 0xCC, 0xCC))
    txt(s5, ticker, Inches(0.65), top + Inches(0.1), Inches(1.2), Inches(0.4),
        size=16, bold=True, color=BURNT_ORANGE)
    txt(s5, body, Inches(0.65), top + Inches(0.45), Inches(7.4), Inches(0.95),
        size=12, color=DARK_GRAY)
    top += Inches(1.6)

# Right: what we delivered + future work
rect(s5, Inches(8.75), Inches(1.2), Inches(4.25), Inches(2.7),
     fill=WHITE, line_color=RGBColor(0xDD, 0xDD, 0xDD))
txt(s5, "WHAT WE DELIVERED", Inches(8.95), Inches(1.3), Inches(4), Inches(0.3),
    size=10, bold=True, color=BURNT_ORANGE)
divider(s5, Inches(1.6), l=Inches(8.95), w=Inches(3.8))
bullets(s5, [
    "All 6 dashboard modules working",
    "33,675 tickers on 5 quant signals",
    "Live transcript fetch \u2014 any U.S. equity",
    "Claude AI synthesis + divergence flag",
], Inches(8.95), Inches(1.7), Inches(3.9), Inches(1.9), size=13)

rect(s5, Inches(8.75), Inches(4.1), Inches(4.25), Inches(2.8),
     fill=WHITE, line_color=RGBColor(0xDD, 0xDD, 0xDD))
txt(s5, "FUTURE WORK", Inches(8.95), Inches(4.2), Inches(4), Inches(0.3),
    size=10, bold=True, color=BURNT_ORANGE)
divider(s5, Inches(4.5), l=Inches(8.95), w=Inches(3.8))
bullets(s5, [
    "Live data feed (remove stale limit)",
    "Backtest composite signal returns",
    "Expand to international equities",
], Inches(8.95), Inches(4.6), Inches(3.9), Inches(1.8), size=13)


# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
out = "presentation.pptx"
prs.save(out)
print(f"Saved -> {out}")
